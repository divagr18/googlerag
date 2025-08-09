# api/core/document_processor.py
import subprocess
import io
import os
import re
import asyncio
from concurrent.futures import ThreadPoolExecutor
import time
from typing import AsyncIterator, List, Tuple, Dict, Optional
from urllib.parse import urlparse
import numpy as np
import fitz  # PyMuPDF
import aiohttp
import hashlib

from .embedding_manager import OptimizedEmbeddingManager
import docx
import pandas as pd
from PIL import Image
import pytesseract
from pptx import Presentation
from .query_expander import DomainQueryExpander
import tempfile
import shutil
from nltk.tokenize import sent_tokenize

# --- CACHING SYSTEM ---
class DocumentCache:
    """Simple in-memory cache for document processing results."""
    
    def __init__(self, max_size: int = 50):
        self._cache: Dict[str, List[Tuple[str, int]]] = {}
        self._access_times: Dict[str, float] = {}
        self._max_size = max_size
        print(f"📋 Document cache initialized (max_size={max_size})")
    
    def _get_cache_key(self, url: str) -> str:
        """Generate a cache key from URL."""
        return hashlib.md5(url.encode('utf-8')).hexdigest()
    
    def _evict_oldest(self):
        """Remove oldest accessed item to make space."""
        if not self._access_times:
            return
        
        oldest_key = min(self._access_times.keys(), key=lambda k: self._access_times[k])
        del self._cache[oldest_key]
        del self._access_times[oldest_key]
        print(f"📋 Cache evicted oldest entry: {oldest_key}")
    
    def get(self, url: str) -> Optional[List[Tuple[str, int]]]:
        """Get cached document data."""
        cache_key = self._get_cache_key(url)
        
        if cache_key in self._cache:
            # Update access time
            self._access_times[cache_key] = time.time()
            print(f"📋 Cache HIT for {url[:50]}...")
            return self._cache[cache_key]
        
        print(f"📋 Cache MISS for {url[:50]}...")
        return None
    
    def set(self, url: str, pages_data: List[Tuple[str, int]]):
        """Cache document data."""
        cache_key = self._get_cache_key(url)
        
        # Make space if needed
        while len(self._cache) >= self._max_size:
            self._evict_oldest()
        
        self._cache[cache_key] = pages_data
        self._access_times[cache_key] = time.time()
        print(f"💾 Cached document data ({len(pages_data)} pages) for {url[:50]}...")
    
    def clear(self):
        """Clear all cached data."""
        self._cache.clear()
        self._access_times.clear()
        print("📋 Cache cleared")
    
    def stats(self) -> Dict[str, int]:
        """Get cache statistics."""
        return {
            "size": len(self._cache),
            "max_size": self._max_size,
            "total_pages_cached": sum(len(pages) for pages in self._cache.values())
        }

# Global cache instance
_document_cache = DocumentCache(max_size=50)

def get_document_cache() -> DocumentCache:
    """Get the global document cache instance."""
    return _document_cache

def smart_word_doc_chunking(text: str, page_num: int, max_chunk_size: int = 1500, overlap_size: int = 150) -> List[Tuple[str, int]]:
    """
    Enhanced chunking specifically for Word documents with dense text and no paragraph breaks.
    Uses sentence-based chunking with overlap for better context preservation.
    """
    if len(text.strip()) < 100:
        return []
    
    try:
        # Use NLTK to split into sentences
        sentences = sent_tokenize(text)
        # Filter out very short sentences
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    except Exception as e:
        print(f"Sentence tokenization failed: {e}. Using fallback method.")
        # Fallback: split by periods and clean up
        sentences = []
        parts = text.split('.')
        for part in parts:
            part = part.strip()
            if len(part) > 20:
                if not part.endswith('.'):
                    part += '.'
                sentences.append(part)
    
    if not sentences:
        # Last resort: just split the text into chunks
        chunks = []
        for i in range(0, len(text), max_chunk_size):
            chunk_text = text[i:i + max_chunk_size].strip()
            if len(chunk_text) > 50:
                chunks.append((chunk_text, page_num))
        return chunks
    
    chunks = []
    current_chunk_sentences = []
    current_chunk_size = 0
    
    for i, sentence in enumerate(sentences):
        sentence_len = len(sentence)
        
        # Check if adding this sentence would exceed chunk size
        if current_chunk_size + sentence_len > max_chunk_size and current_chunk_sentences:
            # Create chunk from current sentences
            chunk_text = ' '.join(current_chunk_sentences)
            if len(chunk_text.strip()) > 100:  # Only add substantial chunks
                chunks.append((chunk_text, page_num))
            
            # Start new chunk with overlap
            # Keep the last 1-2 sentences for context
            overlap_sentences = current_chunk_sentences[-1:] if current_chunk_sentences else []
            overlap_size_current = sum(len(s) for s in overlap_sentences)
            
            # If overlap is too big, reduce it
            while overlap_sentences and overlap_size_current > overlap_size:
                removed = overlap_sentences.pop(0)
                overlap_size_current -= len(removed)
            
            current_chunk_sentences = overlap_sentences + [sentence]
            current_chunk_size = sum(len(s) for s in current_chunk_sentences)
        else:
            # Add sentence to current chunk
            current_chunk_sentences.append(sentence)
            current_chunk_size += sentence_len
    
    # Add the final chunk if it has content
    if current_chunk_sentences:
        chunk_text = ' '.join(current_chunk_sentences)
        if len(chunk_text.strip()) > 100:
            chunks.append((chunk_text, page_num))
    
    print(f"Word doc chunking: Created {len(chunks)} chunks from {len(sentences)} sentences")
    return chunks

async def download_with_aria2c(url: str, destination_path: str):
    """
    Download a file using aria2c if available; otherwise use aiohttp as a fallback.
    """
    print(f"Starting high-speed download for {url}")
    start_time = time.perf_counter()

    command = [
        "aria2c",
        url,
        "--dir", os.path.dirname(destination_path),
        "--out", os.path.basename(destination_path),
        "-x", "8",
        "-s", "8",
        "--quiet=true",
        "--auto-file-renaming=false"
    ]

    def aria2c_exists():
        """Check if aria2c exists in PATH."""
        return shutil.which("aria2c") is not None

    def run_blocking_download():
        try:
            result = subprocess.run(command, capture_output=True, text=True, check=True)
            return result
        except subprocess.CalledProcessError as e:
            raise IOError(f"aria2c download failed with code {e.returncode}: {e.stderr}")

    async def fallback_download():
        """Fallback: async download using aiohttp."""
        async with aiohttp.ClientSession() as session:
            async with session.get(url) as resp:
                resp.raise_for_status()
                with open(destination_path, "wb") as f:
                    async for chunk in resp.content.iter_chunked(1024 * 64):
                        f.write(chunk)

    try:
        if aria2c_exists():
            loop = asyncio.get_running_loop()
            await loop.run_in_executor(None, run_blocking_download)
        else:
            print("aria2c not found, using fallback download method...")
            await fallback_download()

        duration = time.perf_counter() - start_time
        print(f"Download complete in {duration:.2f}s → {destination_path}")

    except Exception as e:
        if os.path.exists(destination_path):
            os.remove(destination_path)
        raise e

async def stream_document(url: str) -> AsyncIterator[bytes]:
    """
    Streams a document by downloading to a temporary location.
    """
    temp_dir = tempfile.mkdtemp()
    original_filename = os.path.basename(urlparse(url).path) or "downloaded_file"
    destination_path = os.path.join(temp_dir, original_filename)

    try:
        await download_with_aria2c(url, destination_path)
        with open(destination_path, "rb") as f:
            content = f.read()
        
        chunk_size = 8192
        for i in range(0, len(content), chunk_size):
            yield content[i:i+chunk_size]

    except Exception as e:
        print(f"FATAL: An error occurred during the download or streaming process: {e}")
        raise
    finally:
        print(f"CLEANUP: Removing temporary directory: {temp_dir}")
        shutil.rmtree(temp_dir, ignore_errors=True)

# --- CACHED DOCUMENT PROCESSING ---
async def process_document_stream_cached(url: str, document_iterator: Optional[AsyncIterator[bytes]] = None, use_cache: bool = True) -> List[Tuple[str, int]]:
    """
    Process document with caching to avoid duplicate processing.
    
    Args:
        url: Document URL
        document_iterator: Optional pre-existing iterator (if None, will create new one)
        use_cache: Whether to use caching (default True)
    
    Returns:
        List of (text, page_num) tuples
    """
    cache = get_document_cache()
    
    # Check cache first (if enabled)
    if use_cache:
        cached_result = cache.get(url)
        if cached_result is not None:
            return cached_result
    
    # If no iterator provided, create one
    if document_iterator is None:
        print(f"📥 Creating new document stream for {url[:50]}...")
        document_iterator = stream_document(url)
    
    # Process the document
    start_time = time.perf_counter()
    pages_data = await process_document_stream(url, document_iterator)
    process_time = time.perf_counter() - start_time
    
    print(f"📄 Document processing took {process_time:.2f}s ({len(pages_data)} pages)")
    
    # Cache the result (if enabled and successful)
    if use_cache and pages_data:
        cache.set(url, pages_data)
    
    return pages_data

cpu_executor = ThreadPoolExecutor(max_workers=os.cpu_count() or 4)

def _extract_page_text_from_bytes(args: Tuple[bytes, int]) -> Tuple[str, int]:
    pdf_bytes, page_no = args
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            return doc.load_page(page_no).get_text(), page_no + 1
    except Exception as e:
        print(f"Error extracting page {page_no}: {e}")
        return "", page_no + 1

async def _extract_text_from_pdf_stream(pdf_stream: io.BytesIO) -> List[Tuple[str, int]]:
    pdf_bytes = pdf_stream.getvalue()
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            page_count = doc.page_count
    except Exception as e:
        raise ValueError(f"Failed to open PDF from memory stream: {e}")
    if page_count == 0:
        return []
    loop = asyncio.get_running_loop()
    tasks = [
        loop.run_in_executor(cpu_executor, _extract_page_text_from_bytes, (pdf_bytes, i))
        for i in range(page_count)
    ]
    return await asyncio.gather(*tasks)

def _extract_text_from_docx_stream(docx_stream: io.BytesIO) -> List[Tuple[str, int]]:
    try:
        document = docx.Document(docx_stream)
        full_text = "\n".join([para.text for para in document.paragraphs])
        
        # Use the enhanced chunking for Word docs
        return smart_paragraph_split(full_text, 1, is_word_doc=True)
    except Exception as e:
        raise ValueError(f"Failed to process DOCX file: {e}")

def _extract_text_from_xlsx_stream(xlsx_stream: io.BytesIO) -> List[Tuple[str, int]]:
    try:
        sheets = pd.read_excel(xlsx_stream, sheet_name=None)
        full_text = [f"Sheet: {name}\n{df.to_csv(index=False)}" for name, df in sheets.items()]
        return [("\n\n".join(full_text), 1)]
    except Exception as e:
        raise ValueError(f"Failed to process XLSX file: {e}")

def _extract_text_from_image_stream(image_stream: io.BytesIO) -> List[Tuple[str, int]]:
    try:
        image = Image.open(image_stream)
        text = pytesseract.image_to_string(image)
        return [(text, 1)]
    except Exception as e:
        raise ValueError(f"Failed to process image file with OCR: {e}")

def _extract_text_from_pptx_stream(pptx_stream: io.BytesIO) -> List[Tuple[str, int]]:
    try:
        presentation = Presentation(pptx_stream)
        pages_data = []
        for i, slide in enumerate(presentation.slides):
            slide_texts = [run.text for shape in slide.shapes if shape.has_text_frame for para in shape.text_frame.paragraphs for run in para.runs]
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_texts.append("\n--- Notes ---\n" + slide.notes_slide.notes_text_frame.text)
            pages_data.append(("\n".join(slide_texts), i + 1))
        return pages_data
    except Exception as e:
        raise ValueError(f"Failed to process PPTX file: {e}")

async def process_document_stream(url: str, document_iterator: AsyncIterator[bytes]) -> List[Tuple[str, int]]:
    path = urlparse(url).path
    file_type = os.path.splitext(path)[1].lower() or ".txt"
    document_bytes_io = io.BytesIO(b"".join([chunk async for chunk in document_iterator]))
    loop = asyncio.get_running_loop()
    
    image_formats = ['.png', '.jpeg', '.jpg', '.bmp', '.tiff', '.gif']
    if file_type == '.pdf': return await _extract_text_from_pdf_stream(document_bytes_io)
    if file_type == '.docx': return await loop.run_in_executor(cpu_executor, _extract_text_from_docx_stream, document_bytes_io)
    if file_type == '.pptx': return await loop.run_in_executor(cpu_executor, _extract_text_from_pptx_stream, document_bytes_io)
    if file_type == '.xlsx': return await loop.run_in_executor(cpu_executor, _extract_text_from_xlsx_stream, document_bytes_io)
    if file_type in image_formats: return await loop.run_in_executor(cpu_executor, _extract_text_from_image_stream, document_bytes_io)
    if file_type in ['.txt', '.md', '.csv']: return [(document_bytes_io.read().decode('utf-8', errors='ignore'), 1)]
    raise ValueError(f"Unsupported file type: '{file_type}'.")

def smart_paragraph_split(text: str, page_num: int, is_word_doc: bool = False) -> List[Tuple[str, int]]:
    """
    Enhanced paragraph splitting with special handling for Word documents.
    """
    # Special handling for Word documents with dense text
    if is_word_doc:
        # Check if this looks like dense text (few line breaks relative to length)
        line_break_ratio = text.count('\n') / len(text) if len(text) > 0 else 0
        if line_break_ratio < 0.001:  # Very few line breaks - likely dense text
            print(f"Detected dense Word document text, using sentence-based chunking")
            return smart_word_doc_chunking(text, page_num)
    
    # Original logic for other documents
    MAX_PARA_LENGTH = 4096
    paragraphs = text.split('\n\n')
    result: List[Tuple[str, int]] = []
    for para in paragraphs:
        para = para.strip()
        if len(para) > 50:
            if len(para) > MAX_PARA_LENGTH:
                for i in range(0, len(para), MAX_PARA_LENGTH):
                    sub_para = para[i:i + MAX_PARA_LENGTH].strip()
                    if len(sub_para) > 50:
                        result.append((sub_para, page_num))
            else:
                result.append((para, page_num))
    return result

async def optimized_semantic_chunk_text(pages_data: List[Tuple[str, int]], embedding_manager: OptimizedEmbeddingManager, similarity_threshold: float = 0.3, max_chunk_size: int = 1500) -> Tuple[List[Dict], np.ndarray]:
    if not pages_data:
        return [], np.array([])
    
    all_paragraphs = []
    for page_text, page_num in pages_data:
        if len(page_text) <= max_chunk_size * 1.2:  # Allow some buffer
            all_paragraphs.append((page_text, page_num))
        else:
            all_paragraphs.extend(smart_paragraph_split(page_text, page_num))
    
    if not all_paragraphs:
        return [], np.array([])
    
    print(f"Starting async-batched chunking on {len(all_paragraphs)} paragraphs...")
    paragraph_texts = [p[0] for p in all_paragraphs]
    
    # Process embeddings in batches
    all_embeddings = []
    for i in range(0, len(paragraph_texts), 256):
        batch = paragraph_texts[i:i+256]
        batch_embeddings = embedding_manager.encode_batch(batch)
        all_embeddings.append(batch_embeddings)
    
    para_embs = np.vstack(all_embeddings)
    
    # Just create chunks directly
    chunks = []
    chunk_embs = []
    
    for i, (text, page_num) in enumerate(all_paragraphs):
        chunks.append({"text": text, "metadata": {"page": page_num}})
        chunk_embs.append(para_embs[i])
    
    print(f"Successfully created {len(chunks)} semantic chunks with pre-computed embeddings.")
    return chunks, np.vstack(chunk_embs)

async def build_enhanced_retrieval_systems(chunks: List[Dict]) -> Optional[DomainQueryExpander]:
    """
    Builds only the DomainQueryExpander. Multi-vector system is removed.
    """
    print("🔧 Building enhanced retrieval systems (Query Expander)...")
    start_time = time.perf_counter()
    
    if not chunks:
        print("⚠️ No chunks provided to build query expander.")
        return None
        
    all_texts = [chunk['text'] for chunk in chunks]
    
    try:
        loop = asyncio.get_running_loop()
        query_expander = await loop.run_in_executor(
            cpu_executor,
            lambda: DomainQueryExpander(all_texts, min_term_freq=2, max_expansion_terms=4)
        )
        print(f"✅ Query Expander built in {time.perf_counter() - start_time:.2f}s")
        return query_expander
    except Exception as e:
        print(f"⚠️ Query expander build failed: {e}")
        return None

# --- UTILITY FUNCTIONS ---
def clear_document_cache():
    """Clear the document cache (useful for testing or memory management)."""
    get_document_cache().clear()

def get_cache_stats() -> Dict[str, int]:
    """Get document cache statistics."""
    return get_document_cache().stats()